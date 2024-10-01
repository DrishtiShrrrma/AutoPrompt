from langchain_core.tools import tool
from typing import Optional
from .browser import SimpleTextBrowser
import requests
import re
import mimetypes
from markdownify import MarkdownConverter
from youtube_transcript_api import YouTubeTranscriptApi
import pdfminer.high_level
import mammoth
import pptx
from bs4 import BeautifulSoup

# Initialize the browser instance
browser = SimpleTextBrowser()

# Helper function to retrieve browser state
def _browser_state() -> str:
    header = f"Address: {browser.address}\n"
    if browser.page_title is not None:
        header += f"Title: {browser.page_title}\n"
    
    current_page = browser.viewport_current_page
    total_pages = len(browser.viewport_pages)
    header += f"Viewport position: Showing page {current_page+1} of {total_pages}."
    return header

@tool(name="web_search", description="Perform a web search query and return the search results.")
def web_search_tool(query: str, filter_year: Optional[int] = None) -> str:
    """Perform a web search query and return the search results."""
    browser.visit_page(f"google: {query}", filter_year=filter_year)
    header = _browser_state()
    return header + "\n=======================\n" + browser.viewport

@tool(name="navigational_web_search", description="Perform a navigational web search and visit the top result.")
def navigational_web_search_tool(query: str) -> str:
    """Perform a navigational web search and visit the top result."""
    browser.visit_page(f"google: {query}")
    
    # Extract the first link and navigate to it
    match = re.search(r"\[.*?\]\((http.*?)\)", browser.page_content)
    if match:
        browser.visit_page(match.group(1))
    
    header = _browser_state()
    return header + "\n=======================\n" + browser.viewport

@tool(name="visit_page", description="Visit a webpage at a given URL and return its text.")
def visit_page_tool(url: str) -> str:
    """Visit a webpage at a given URL and return its text."""
    browser.visit_page(url)
    header = _browser_state()
    return header + "\n=======================\n" + browser.viewport

@tool(name="download_file", description="Download a file at a given URL (non-pdf/txt/html).")
def download_file_tool(url: str) -> str:
    """Download a file at a given URL (use visit_page for PDF/TXT/HTML files)."""
    if "arxiv" in url:
        url = url.replace("abs", "pdf")
    response = requests.get(url)
    content_type = response.headers.get("content-type", "")
    extension = mimetypes.guess_extension(content_type)
    
    if extension:
        file_path = f"./downloads/file{extension}"
    else:
        file_path = "./downloads/file.object"
    
    with open(file_path, "wb") as f:
        f.write(response.content)

    if "pdf" in extension or "txt" in extension or "htm":
        raise Exception("Do not use this tool for PDF/TXT/HTML files, use visit_page instead.")
    
    return f"File was downloaded and saved under path {file_path}."

@tool(name="page_up", description="Scroll the viewport UP one page-length in the current webpage.")
def page_up_tool() -> str:
    """Scroll the viewport UP one page-length."""
    browser.page_up()
    header = _browser_state()
    return header + "\n=======================\n" + browser.viewport

@tool(name="page_down", description="Scroll the viewport DOWN one page-length in the current webpage.")
def page_down_tool() -> str:
    """Scroll the viewport DOWN one page-length."""
    browser.page_down()
    header = _browser_state()
    return header + "\n=======================\n" + browser.viewport

@tool(name="find_on_page", description="Find the first occurrence of a string on the current page (Ctrl+F).")
def find_on_page_tool(search_string: str) -> str:
    """Find the first occurrence of a string on the current page (Ctrl+F)."""
    find_result = browser.find_on_page(search_string)
    header = _browser_state()
    
    if find_result is None:
        return header + f"\n=======================\nThe search string '{search_string}' was not found on this page."
    else:
        return header + "\n=======================\n" + browser.viewport

@tool(name="find_next", description="Find the next occurrence of the previously searched string.")
def find_next_tool() -> str:
    """Find the next occurrence of the previously searched string."""
    find_result = browser.find_next()
    header = _browser_state()
    
    if find_result is None:
        return header + "\n=======================\nThe search string was not found on this page."
    else:
        return header + "\n=======================\n" + browser.viewport

@tool(name="find_archived_url", description="Search Wayback Machine for the archived version of a URL for a given date.")
def find_archived_url_tool(url: str, date: str) -> str:
    """Search Wayback Machine for the archived version of a URL for a given date."""
    archive_url = f"https://archive.org/wayback/available?url={url}&timestamp={date}"
    response = requests.get(archive_url).json()
    try:
        closest = response["archived_snapshots"]["closest"]
    except KeyError:
        raise Exception(f"URL was not archived on Wayback Machine for the given date.")
    
    browser.visit_page(closest["url"])
    header = _browser_state()
    return f"Web archive for URL {url}, snapshot taken at date {closest['timestamp'][:8]}:\n" + header + "\n=======================\n" + browser.viewport

# Newly added tools from mdconvert.py

@tool(name="plain_text_conversion", description="Convert a plain text file into a text string.")
def plain_text_conversion_tool(file_path: str, file_extension: Optional[str] = None) -> str:
    """Convert a plain text file into a text string."""
    if file_extension == "":
        return "Error: Invalid file extension."

    content_type, encoding = mimetypes.guess_type(f"__placeholder{file_extension}")

    with open(file_path, "rt") as file_handle:
        text_content = file_handle.read()

    return text_content

@tool(name="html_conversion", description="Convert an HTML file into markdown text.")
def html_conversion_tool(file_path: str, file_extension: Optional[str] = None) -> str:
    """Convert an HTML file into markdown text."""
    if file_extension.lower() not in [".html", ".htm"]:
        return "Error: Invalid file extension for HTML."

    with open(file_path, "rt") as file_handle:
        html_content = file_handle.read()

    soup = BeautifulSoup(html_content, "html.parser")
    
    # Remove script and style blocks
    for script in soup(["script", "style"]):
        script.extract()

    body_elm = soup.find("body")
    markdown_text = ""
    if body_elm:
        markdown_text = MarkdownConverter().convert_soup(body_elm)
    else:
        markdown_text = MarkdownConverter().convert_soup(soup)

    return markdown_text

@tool(name="youtube_transcript", description="Extract the transcript of a YouTube video using the video URL.")
def youtube_transcript_tool(url: str) -> str:
    """Extract the transcript of a YouTube video using the video URL."""
    parsed_url = urlparse(url)
    params = parse_qs(parsed_url.query)

    if "v" not in params:
        return "Error: Invalid YouTube URL."

    video_id = params["v"][0]
    transcript = YouTubeTranscriptApi.get_transcript(video_id)

    transcript_text = " ".join([part["text"] for part in transcript])

    return transcript_text

@tool(name="pdf_conversion", description="Convert a PDF file into a text string.")
def pdf_conversion_tool(file_path: str, file_extension: Optional[str] = None) -> str:
    """Convert a PDF file into a text string."""
    if file_extension.lower() != ".pdf":
        return "Error: Invalid file extension for PDF."

    text_content = pdfminer.high_level.extract_text(file_path)
    return text_content

@tool(name="docx_conversion", description="Convert a DOCX file into a markdown text.")
def docx_conversion_tool(file_path: str, file_extension: Optional[str] = None) -> str:
    """Convert a DOCX file into a markdown text."""
    if file_extension.lower() != ".docx":
        return "Error: Invalid file extension for DOCX."

    with open(file_path, "rb") as docx_file:
        result = mammoth.convert_to_html(docx_file)
        html_content = result.value

    return html_content

@tool(name="pptx_conversion", description="Convert a PPTX file into a markdown text.")
def pptx_conversion_tool(file_path: str, file_extension: Optional[str] = None) -> str:
    """Convert a PPTX file into a markdown text."""
    if file_extension.lower() != ".pptx":
        return "Error: Invalid file extension for PPTX."

    md_content = ""
    presentation = pptx.Presentation(file_path)

    for slide in presentation.slides:
        title = slide.shapes.title
        md_content += f"# {title.text}\n"
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                md_content += f"{shape.text}\n"

    return md_content
